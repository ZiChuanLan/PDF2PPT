from __future__ import annotations

import re
import statistics
from pathlib import Path

import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Pt

from .geometry import px_bbox_to_pt_bbox
from .models import OcrLine, PageResult, VisualRegion

_EMU_PER_PT = 12700.0


def _pt_to_emu(value_pt: float) -> int:
    return int(round(float(value_pt) * _EMU_PER_PT))


def _contains_cjk(text: str) -> bool:
    for ch in text:
        code = ord(ch)
        if (
            0x4E00 <= code <= 0x9FFF
            or 0x3400 <= code <= 0x4DBF
            or 0xF900 <= code <= 0xFAFF
        ):
            return True
    return False


def _estimate_font_size_pt(*, text: str, box_w_pt: float, box_h_pt: float) -> float:
    compact = [ch for ch in str(text or "") if not ch.isspace()]
    if not compact:
        return 10.0

    length = len(compact)
    cjk_count = sum(1 for ch in compact if _contains_cjk(ch))
    cjk_ratio = float(cjk_count) / float(max(1, length))

    if cjk_ratio >= 0.7:
        width_factor = 0.95
    elif cjk_ratio <= 0.05:
        width_factor = 0.56
    else:
        width_factor = 0.72

    width_limit = float(box_w_pt) / max(1.0, float(length) * width_factor)
    height_limit = float(box_h_pt) * 0.72
    estimated = min(width_limit, height_limit)
    if length <= 3:
        estimated = min(estimated, height_limit * 0.9)
    return max(7.0, min(34.0, estimated))


def _estimate_font_size_pt_height_first(*, text: str, box_w_pt: float, box_h_pt: float) -> float:
    compact = [ch for ch in str(text or "") if not ch.isspace()]
    if not compact:
        return 10.0
    length = len(compact)
    has_cjk = _contains_cjk(text)

    # Height-first for better visual parity on scanned slides.
    height_limit = float(box_h_pt) * 0.82

    if has_cjk:
        width_factor = 1.05
    else:
        width_factor = 0.62

    width_limit = float(box_w_pt) / max(1.0, float(length) * width_factor)
    estimated = min(width_limit, height_limit)
    if len(text.strip()) <= 3:
        estimated = min(estimated, height_limit * 0.9)
    return max(8.0, min(42.0, estimated))


def _is_low_value_noise(line: OcrLine) -> bool:
    text = str(line.text or "").strip()
    if not text:
        return True
    compact = [ch for ch in text if not ch.isspace()]
    if not compact:
        return True

    if len(compact) == 1:
        return not _contains_cjk(text)

    noise_pattern = re.search(r"[A-Za-z0-9][|_=]{1,}[A-Za-z0-9]", text) is not None
    ascii_chars = sum(1 for ch in compact if ord(ch) < 128)
    ascii_ratio = float(ascii_chars) / float(len(compact))

    if noise_pattern:
        return True
    if _contains_cjk(text) and ascii_ratio >= 0.65 and len(compact) >= 8:
        return True

    return False


def _iter_clean_lines(lines: list[OcrLine]) -> list[OcrLine]:
    cleaned: list[OcrLine] = []
    seen: set[tuple[str, int, int, int, int]] = set()
    for line in lines:
        if _is_low_value_noise(line):
            continue
        x0, y0, x1, y1 = [float(v) for v in line.bbox]
        key = (
            line.text.strip(),
            int(round(x0)),
            int(round(y0)),
            int(round(x1)),
            int(round(y1)),
        )
        if key in seen:
            continue
        seen.add(key)
        cleaned.append(line)

    cleaned.sort(key=lambda item: (float(item.bbox[1]), float(item.bbox[0])))
    return cleaned


def _iter_clean_regions(regions: list[VisualRegion], *, width_px: int, height_px: int) -> list[VisualRegion]:
    page_area = float(max(1, width_px * height_px))
    cleaned: list[VisualRegion] = []

    for region in regions:
        x0, y0, x1, y1 = [float(v) for v in region.bbox]
        box_w = max(1.0, x1 - x0)
        box_h = max(1.0, y1 - y0)
        area_ratio = (box_w * box_h) / page_area
        if area_ratio < 0.0025:
            continue
        if area_ratio > 0.82:
            continue
        if box_h / max(1.0, float(height_px)) < 0.03:
            continue
        cleaned.append(region)

    cleaned.sort(key=lambda item: (float(item.bbox[1]), float(item.bbox[0])))
    return cleaned


def _boxes_overlap(a: list[float], b: list[float]) -> bool:
    ax0, ay0, ax1, ay1 = [float(v) for v in a]
    bx0, by0, bx1, by1 = [float(v) for v in b]
    ix0 = max(ax0, bx0)
    iy0 = max(ay0, by0)
    ix1 = min(ax1, bx1)
    iy1 = min(ay1, by1)
    if ix1 <= ix0 or iy1 <= iy0:
        return False
    inter = (ix1 - ix0) * (iy1 - iy0)
    area_a = max(1.0, (ax1 - ax0) * (ay1 - ay0))
    area_b = max(1.0, (bx1 - bx0) * (by1 - by0))
    cover_a = inter / area_a
    cover_b = inter / area_b
    return cover_a >= 0.42 or cover_b >= 0.42


def _exclude_lines_inside_regions(
    lines: list[OcrLine],
    regions: list[VisualRegion],
) -> list[OcrLine]:
    if not lines or not regions:
        return lines
    kept: list[OcrLine] = []
    region_bboxes = [[float(v) for v in region.bbox] for region in regions]
    for line in lines:
        line_bbox = [float(v) for v in line.bbox]
        if any(_boxes_overlap(line_bbox, region_bbox) for region_bbox in region_bboxes):
            continue
        kept.append(line)
    kept.sort(key=lambda item: (float(item.bbox[1]), float(item.bbox[0])))
    return kept


def _group_lines(lines: list[OcrLine]) -> list[list[OcrLine]]:
    if not lines:
        return []

    buckets: dict[str, list[OcrLine]] = {}
    order: list[str] = []
    for idx, line in enumerate(lines):
        if line.group_id:
            key = line.group_id
        else:
            key = f"line:{idx}"
        if key not in buckets:
            buckets[key] = []
            order.append(key)
        buckets[key].append(line)

    groups = [buckets[key] for key in order]
    groups.sort(
        key=lambda group: (
            float(min(item.bbox[1] for item in group)),
            float(min(item.bbox[0] for item in group)),
        )
    )
    return groups


def _merge_bbox(lines: list[OcrLine]) -> list[float]:
    x0 = min(float(line.bbox[0]) for line in lines)
    y0 = min(float(line.bbox[1]) for line in lines)
    x1 = max(float(line.bbox[2]) for line in lines)
    y1 = max(float(line.bbox[3]) for line in lines)
    return [x0, y0, x1, y1]


def _sample_text_color(
    *,
    image: Image.Image | None,
    line_bbox_px: list[float],
    width_px: int,
    height_px: int,
) -> RGBColor:
    if image is None:
        return RGBColor(0, 0, 0)

    x0, y0, x1, y1 = [float(v) for v in line_bbox_px]
    ix0 = max(0, int(round(x0)))
    iy0 = max(0, int(round(y0)))
    ix1 = min(width_px, int(round(x1)))
    iy1 = min(height_px, int(round(y1)))
    if ix1 <= ix0 or iy1 <= iy0:
        return RGBColor(0, 0, 0)

    crop = image.crop((ix0, iy0, ix1, iy1)).convert("RGB")
    arr = np.array(crop)
    if arr.ndim != 3 or arr.shape[2] < 3:
        return RGBColor(0, 0, 0)
    flat = arr.reshape(-1, arr.shape[2])
    rgb_pixels = [
        (int(px[0]), int(px[1]), int(px[2]))
        for px in flat
    ]
    if not rgb_pixels:
        return RGBColor(0, 0, 0)

    # Use dark-percentile to approximate actual glyph color.
    lum_list: list[tuple[float, tuple[int, int, int]]] = []
    for r, g, b in rgb_pixels:
        lum = 0.299 * float(r) + 0.587 * float(g) + 0.114 * float(b)
        lum_list.append((lum, (int(r), int(g), int(b))))

    lum_list.sort(key=lambda item: item[0])
    take = max(20, int(len(lum_list) * 0.1))
    selected = lum_list[:take]
    sr = int(round(sum(px[0] for _, px in selected) / float(len(selected))))
    sg = int(round(sum(px[1] for _, px in selected) / float(len(selected))))
    sb = int(round(sum(px[2] for _, px in selected) / float(len(selected))))

    # Clamp too-light colors toward dark to preserve readability.
    mean_lum = 0.299 * sr + 0.587 * sg + 0.114 * sb
    if mean_lum > 165:
        scale = 165.0 / max(1.0, mean_lum)
        sr = int(round(sr * scale))
        sg = int(round(sg * scale))
        sb = int(round(sb * scale))

    return RGBColor(max(0, min(255, sr)), max(0, min(255, sg)), max(0, min(255, sb)))


def build_ppt_from_pages(page_results: list[PageResult], out_path: Path) -> Path:
    if not page_results:
        raise ValueError("No pages to build PPT")

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    first = page_results[0].render
    prs.slide_width = Emu(_pt_to_emu(first.width_pt))
    prs.slide_height = Emu(_pt_to_emu(first.height_pt))

    for page in page_results:
        slide = prs.slides.add_slide(blank_layout)
        page_image_for_color: Image.Image | None = None
        try:
            page_image_for_color = Image.open(page.render.image_path).convert("RGB")
        except Exception:
            page_image_for_color = None

        page_w_pt = page.render.width_pt
        page_h_pt = page.render.height_pt

        clean_regions = _iter_clean_regions(
            page.image_regions,
            width_px=page.render.width_px,
            height_px=page.render.height_px,
        )
        base_lines = _iter_clean_lines(page.ocr_lines)
        clean_lines = _exclude_lines_inside_regions(base_lines, clean_regions)
        if base_lines and clean_regions:
            removed_ratio = 1.0 - (float(len(clean_lines)) / float(len(base_lines)))
            if removed_ratio >= 0.8:
                clean_lines = base_lines

        slide.shapes.add_picture(
            str(page.cleaned_image_path),
            Emu(0),
            Emu(0),
            Emu(_pt_to_emu(page_w_pt)),
            Emu(_pt_to_emu(page_h_pt)),
        )

        for region in clean_regions:
            if region.crop_path is None or not Path(region.crop_path).exists():
                continue
            bbox_pt = px_bbox_to_pt_bbox(
                region.bbox,
                image_width_px=page.render.width_px,
                image_height_px=page.render.height_px,
                page_width_pt=page_w_pt,
                page_height_pt=page_h_pt,
            )
            x0, y0, x1, y1 = bbox_pt
            w = max(1.0, x1 - x0)
            h = max(1.0, y1 - y0)
            slide.shapes.add_picture(
                str(region.crop_path),
                Emu(_pt_to_emu(x0)),
                Emu(_pt_to_emu(y0)),
                Emu(_pt_to_emu(w)),
                Emu(_pt_to_emu(h)),
            )

        grouped_lines = _group_lines(clean_lines)

        for group in grouped_lines:
            if not group:
                continue

            line_layouts: list[tuple[OcrLine, tuple[float, float, float, float], float, float]] = []
            line_sizes: list[float] = []
            line_colors: list[RGBColor] = []

            for line in group:
                bbox_pt = px_bbox_to_pt_bbox(
                    line.bbox,
                    image_width_px=page.render.width_px,
                    image_height_px=page.render.height_px,
                    page_width_pt=page_w_pt,
                    page_height_pt=page_h_pt,
                )
                x0, y0, x1, y1 = bbox_pt
                w = max(1.0, x1 - x0)
                h = max(1.0, y1 - y0)
                line_layouts.append((line, (x0, y0, x1, y1), w, h))

                line_sizes.append(
                    _estimate_font_size_pt_height_first(
                        text=line.text,
                        box_w_pt=w,
                        box_h_pt=h,
                    )
                )
                line_colors.append(
                    _sample_text_color(
                        image=page_image_for_color,
                        line_bbox_px=line.bbox,
                        width_px=page.render.width_px,
                        height_px=page.render.height_px,
                    )
                )

            if line_sizes:
                group_font_size = float(statistics.median(line_sizes))
            else:
                group_font_size = 12.0
            group_font_size = max(10.0, min(42.0, group_font_size))

            if line_colors:
                def color_lum(color: RGBColor) -> float:
                    return 0.299 * float(color[0]) + 0.587 * float(color[1]) + 0.114 * float(color[2])

                group_color = min(line_colors, key=color_lum)
            else:
                group_color = RGBColor(0, 0, 0)

            for line, (x0, y0, x1, y1), w, h in line_layouts:

                tx = slide.shapes.add_textbox(
                    Emu(_pt_to_emu(x0)),
                    Emu(_pt_to_emu(y0)),
                    Emu(_pt_to_emu(w)),
                    Emu(_pt_to_emu(h)),
                )
                tf = tx.text_frame
                tf.clear()
                tf.word_wrap = False
                tf.auto_size = MSO_AUTO_SIZE.NONE
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = line.text
                run.font.size = Pt(group_font_size)
                run.font.bold = False
                run.font.color.rgb = group_color

        if page_image_for_color is not None:
            try:
                page_image_for_color.close()
            except Exception:
                pass

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path
