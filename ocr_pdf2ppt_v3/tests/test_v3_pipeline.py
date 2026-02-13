import json
from pathlib import Path

import fitz
from pptx import Presentation

from ocr_pdf2ppt_v3.models import OcrLine, VisualRegion
from ocr_pdf2ppt_v3.pipeline import convert_pdf_to_ppt


class FakeOcrClient:
    def __init__(self, *, api_key: str, base_url: str, model: str, max_image_side_px: int = 2200):
        self.api_key = api_key
        self.base_url = base_url
        self.model = model
        self.max_image_side_px = max_image_side_px

    def ocr_page(self, page, *, pass_mode: str = "primary"):
        sx = page.width_px / page.width_pt
        sy = page.height_px / page.height_pt
        return [
            OcrLine(
                text="Hello OCR V3",
                bbox=[80 * sx, 60 * sy, 260 * sx, 95 * sy],
                confidence=0.98,
                source="ai_primary" if pass_mode == "primary" else "ai_retry",
            )
        ]

    def detect_layout_regions(self, page):
        sx = page.width_px / page.width_pt
        sy = page.height_px / page.height_pt
        return [
            VisualRegion(
                bbox=[300 * sx, 200 * sy, 500 * sx, 430 * sy],
                label="chart",
                confidence=0.92,
                source="ai_layout",
            )
        ]


def _make_sample_pdf(path: Path) -> None:
    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    page.draw_rect(fitz.Rect(220, 220, 500, 450), color=(0.1, 0.2, 0.9), fill=(0.3, 0.7, 0.9))
    page.insert_text((80, 90), "Hello OCR V3", fontsize=24, color=(0, 0, 0))
    doc.save(str(path))
    doc.close()


def test_convert_pdf_to_ppt_end_to_end(tmp_path: Path, monkeypatch) -> None:
    input_pdf = tmp_path / "input.pdf"
    output_pptx = tmp_path / "output.pptx"
    work_dir = tmp_path / "work"
    _make_sample_pdf(input_pdf)

    import ocr_pdf2ppt_v3.pipeline as pipeline

    monkeypatch.setattr(pipeline, "SiliconFlowOcrClient", FakeOcrClient)

    result = convert_pdf_to_ppt(
        input_pdf=input_pdf,
        output_pptx=output_pptx,
        api_key="dummy-key",
        base_url="https://api.siliconflow.cn/v1",
        model="dummy-model",
        max_pages=1,
        work_dir=work_dir,
    )

    assert result.output_pptx.exists()
    assert result.pages == 1
    assert (work_dir / "pages" / "page-0001.png").exists()
    assert (work_dir / "clean" / "page-0001.clean.png").exists()
    assert (work_dir / "ocr" / "page-0001.json").exists()
    assert (work_dir / "ocr" / "page-0001.layout.json").exists()
    assert result.debug_json.exists()

    prs = Presentation(str(output_pptx))
    assert len(prs.slides) == 1
    shapes = prs.slides[0].shapes
    assert len(shapes) >= 3  # background + image + text


def test_convert_pdf_to_ppt_ai_retry(tmp_path: Path, monkeypatch) -> None:
    input_pdf = tmp_path / "input.pdf"
    output_pptx = tmp_path / "output.pptx"
    work_dir = tmp_path / "work"
    _make_sample_pdf(input_pdf)

    import ocr_pdf2ppt_v3.pipeline as pipeline

    class TwoPassClient(FakeOcrClient):
        def ocr_page(self, page, *, pass_mode: str = "primary"):
            if pass_mode == "primary":
                return []
            return [
                OcrLine(
                    text="Retry Text",
                    bbox=[80, 60, 260, 100],
                    confidence=0.99,
                    source="ai_retry",
                )
            ]

    monkeypatch.setattr(pipeline, "SiliconFlowOcrClient", TwoPassClient)

    result = convert_pdf_to_ppt(
        input_pdf=input_pdf,
        output_pptx=output_pptx,
        api_key="dummy-key",
        base_url="https://api.siliconflow.cn/v1",
        model="dummy-model",
        max_pages=1,
        work_dir=work_dir,
    )

    assert result.fallback_pages == 1
    assert result.empty_pages == 0


def test_convert_pdf_to_ppt_auto_backend_paddle(tmp_path: Path, monkeypatch) -> None:
    input_pdf = tmp_path / "input.pdf"
    output_pptx = tmp_path / "output.pptx"
    work_dir = tmp_path / "work"
    _make_sample_pdf(input_pdf)

    import ocr_pdf2ppt_v3.pipeline as pipeline

    class FakePaddleClient:
        def __init__(self, *, api_key: str, base_url: str, model: str, max_image_side_px: int = 2200):
            self.api_key = api_key
            self.base_url = base_url
            self.model = model
            self.max_image_side_px = max_image_side_px

        def ocr_page(self, page, *, pass_mode: str = "primary"):
            sx = page.width_px / page.width_pt
            sy = page.height_px / page.height_pt
            return [
                OcrLine(
                    text="Paddle Structured",
                    bbox=[90 * sx, 100 * sy, 320 * sx, 145 * sy],
                    confidence=0.99,
                    source="ai_primary" if pass_mode == "primary" else "ai_retry",
                )
            ]

        def detect_layout_regions(self, page):
            return []

    class ExplodingOpenAIClient:
        def __init__(self, **kwargs):
            raise AssertionError("openai chat backend should not be selected for Paddle model")

    monkeypatch.setattr(pipeline, "PaddleDocParserOcrClient", FakePaddleClient)
    monkeypatch.setattr(pipeline, "SiliconFlowOcrClient", ExplodingOpenAIClient)

    result = convert_pdf_to_ppt(
        input_pdf=input_pdf,
        output_pptx=output_pptx,
        api_key="dummy-key",
        base_url="https://api.siliconflow.cn/v1",
        model="PaddlePaddle/PaddleOCR-VL",
        max_pages=1,
        work_dir=work_dir,
    )

    assert result.output_pptx.exists()
    debug_payload = json.loads(result.debug_json.read_text(encoding="utf-8"))
    assert debug_payload["ocr_backend"] == "paddle_doc_parser"
