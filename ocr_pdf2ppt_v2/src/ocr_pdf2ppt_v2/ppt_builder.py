from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Pt

from .geometry import px_bbox_to_pt_bbox
from .models import PageResult

_EMU_PER_PT = 12700.0


def _pt_to_emu(value_pt: float) -> int:
    return int(round(float(value_pt) * _EMU_PER_PT))


def build_ppt_from_pages(page_results: list[PageResult], out_path: Path) -> Path:
    if not page_results:
        raise ValueError("No pages to build PPT")

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]

    first = page_results[0].render
    prs.slide_width = Emu(_pt_to_emu(first.width_pt))
    prs.slide_height = Emu(_pt_to_emu(first.height_pt))

    for page in page_results:
        slide = prs.slides.add_slide(blank_layout)

        page_w_pt = page.render.width_pt
        page_h_pt = page.render.height_pt

        slide.shapes.add_picture(
            str(page.cleaned_image_path),
            Emu(0),
            Emu(0),
            Emu(_pt_to_emu(page_w_pt)),
            Emu(_pt_to_emu(page_h_pt)),
        )

        for line in page.ocr_lines:
            bbox_pt = px_bbox_to_pt_bbox(
                line.bbox,
                image_width_px=page.render.width_px,
                image_height_px=page.render.height_px,
                page_width_pt=page_w_pt,
                page_height_pt=page_h_pt,
            )
            x0, y0, x1, y1 = bbox_pt
            w = max(1.0, x1 - x0)
            h = max(1.0, y1 - y0)

            tx = slide.shapes.add_textbox(
                Emu(_pt_to_emu(x0)),
                Emu(_pt_to_emu(y0)),
                Emu(_pt_to_emu(w)),
                Emu(_pt_to_emu(h)),
            )
            tf = tx.text_frame
            tf.clear()
            tf.word_wrap = False
            tf.auto_size = MSO_AUTO_SIZE.NONE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = line.text
            run.font.size = Pt(max(7.0, min(64.0, h * 0.75)))
            run.font.bold = False
            run.font.color.rgb = RGBColor(0, 0, 0)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path
